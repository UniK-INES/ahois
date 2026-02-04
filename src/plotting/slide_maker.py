"""
Automates the creation of PowerPoint presentations from simulation result plots.

This module provides the `Slide_maker` class, a tool designed to automatically
assemble generated plots from the `Scenario_comparator` into a structured
PowerPoint presentation. It is driven by an Excel configuration file that
defines different presentation variants, specifying which scenarios to include.

For each variant, the script first ensures the necessary comparative plots are
generated. It then programmatically constructs a `.pptx` file using a template,
adding title slides, section dividers, and individual slides for each plot.
This enables the rapid generation reports directly from the model's output.

:Authors:
 - Ivan Digel <ivan.digel@uni-kassel.de>
 - Sascha Holzhauer <sascha.holzhauer@uni-kassel.de>
 """
import os
import re
import io
import time
import yaml
import matplotlib.pyplot as plt
import pandas as pd
import dill as pickle
import logging
import copy
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches, Cm, Pt
from helpers.config import settings, get_output_path, config_logging
from helpers.i18n import _
from plotting.Scenario_comparator import Scenario_comparator
import uuid
from pptx.oxml import parse_xml
from pptx.oxml.xmlchemy import OxmlElement

logger = logging.getLogger("ahoi.slides")

namespaceURL = {
    "mc":  "http://schemas.openxmlformats.org/markup-compatibility/2006",
    "p":   "http://schemas.openxmlformats.org/presentationml/2006/main",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    "p15": "http://schemas.microsoft.com/office/powerpoint/2012/main",
    "a":   "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r":   "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
}

def namespacesFragment(prefixes):
    """
    Creates an XML namespace declaration string from a list of prefixes.

    This is a low-level helper function for the XML manipulation required to
    create PowerPoint sections. It takes a list of namespace prefixes (e.g.,
    'p', 'a', 'r') and constructs a string of `xmlns` attributes that can be
    inserted into an XML tag.

    Parameters
    ----------
    prefixes : list[str]
        A list of namespace prefixes to include in the declaration.

    Returns
    -------
    str
        A string containing the formatted XML namespace attributes.
    """
    xml = ""
    for prefix in prefixes:
        xml += 'xmlns:' + prefix + '="' + namespaceURL[prefix] +'" '

    return xml
    
class Slide_maker():
    """
    Creates PowerPoint slide decks from Scenario_comparator figures.
    """
    
    def __init__(self):
        """
        Initializes the Slide_maker.

        This constructor sets up the initial state, including the run ID and
        the output path for the generated presentations. It prepares the
        instance variables that will be populated by reading the Excel
        configuration file.
        
        Input:
            column from  "plotting/AHOI_SlideGeneration.xlsx" 
        """
        self.run_id = settings.main.run_id
        self.scenarios = []
        self.prefixes = []
        self.plots_data = None
        self.input_table = None
        self.output_path = get_output_path(runid=self.run_id, subfolder='slides')
        

    def perform_all_processing(self):
        """
        Handles the entire slide generation process.

        This is the main public method that drives the creation of the
        presentations. It reads the Excel configuration file to identify which
        presentation variants to build. For each variant, it calls other
        methods to:
        1.  Configure the scenarios and prefixes.
        2.  Generate the required comparative plots using `Scenario_comparator`.
        3.  Assemble the plots into a new PowerPoint presentation.
        4.  Save the final `.pptx` file.
        """
        file_path = f"{get_output_path(runid=settings.main.run_id, subfolder='slides')}/{settings.slides.config}"
        config_df = pd.read_excel(file_path, sheet_name="Sheet1", nrows=11, index_col=0).transpose()
        config_columns = [col for _,col in config_df.iterrows() if col.active==1]

        logger.debug(f"Columns to handle: {config_columns}")
        prs = Presentation(settings.slides.template)
        self.add_section_start(prs)
        
        colindex = 0
        for column in config_columns:
            colindex +=1
            self.configuration = column
            self.read_input_table()
            settings.scenario_comparison.run_ids = [self.run_id]
            settings.scenario_comparison.scenarios = self.scenarios
            settings.scenario_comparison.files_prefixes = self.prefixes
            self.folder_suffix = "-".join(self.prefixes)
            self.main_img_folder = f"{get_output_path(runid=settings.main.run_id, subfolder='plots')}/{self.folder_suffix}"
            self.attribute_img_folder = f"{self.main_img_folder}/attribute_comparison"
        

            self.make_plots()
            self.add_title_slide(prs)
            if settings.slides.sections:
                self.add_section(prs)
            self.move_slide(prs, 1,len(prs.slides))
            self.add_slides(prs)
            
            if column["split"] == 1 or colindex == len(config_columns):
                # Remove template slide
                xml_slides = prs.slides._sldIdLst
                slides = list(xml_slides)
                slide_id = slides[0]
                xml_slides.remove(slide_id)
                if settings.slides.sections:
                    self.add_section_end(prs)
                
                # Save presentation
                output_file = os.path.join(self.output_path, f"{self.configuration.name}.pptx")
                prs.save(output_file)
                logger.info(f"Slides were saved as {output_file}")
        
                prs = Presentation(settings.slides.template)
                if settings.slides.sections:
                    self.add_section_start(prs)
    
    def make_plots(self):
        """
        Generates the necessary comparative plots for the presentation.

        This method ensures that all required plot images are available before
        slide creation begins. It instantiates the `Scenario_comparator` with
        the appropriate scenarios and runs for the current presentation
        variant. To improve performance, it utilises a caching system, loading
        a pre-processed `comparator` object from a pickle file if one exists,
        or creating and saving one if it does not.
        """
        logger.info("Running scenario comparison for slide maker")
        start_time = time.time()
        picklefile = f"{get_output_path(runid=self.run_id, subfolder='postprocessed')}/" + \
                     f"scenario_comparator_{'-'.join(settings.scenario_comparison.files_prefixes)}_" + \
                     f"{'-'.join([str(i) for i in settings.scenario_comparison.run_ids])}.pkl"
        logger.info(f"Requested pickle file: {picklefile}")
        if not os.path.exists(picklefile):
            comparator = Scenario_comparator(
                        scenarios=self.scenarios,
                        run_ids=settings.scenario_comparison.run_ids,
                        files_prefixes=self.prefixes,
                    )
            with open(picklefile, 'wb') as outp:
                pickle.dump(comparator, outp, pickle.HIGHEST_PROTOCOL)
            
        else:
            with open(picklefile, 'rb') as inp:
                comparator = pickle.load(inp)
            with open(settings.data.plt_settings, "r") as configfile:
                config = yaml.safe_load(configfile)
                
            if "Layout" in config:
                plt.rcParams.update(config["Layout"])
            else:
                raise ValueError("Invalid plotting configuration file: 'Layout' section missing.")
        
        end_time = time.time()
        elapsed_time = end_time - start_time
        logger.info(f"Scenario_comparator initialization took {elapsed_time:.2f} seconds")
        
        start_time = time.time()
        
        # TODO get information from xlsx-sheet which figures to produce (again)
        settings.eval.compare_fulfillment = False #True
        settings.eval.compare_emissions = False #True
        settings.eval.compare_heating_systems_distribution = False #True
        settings.eval.compare_hs_knowledge = False #True
        settings.eval.compare_obstacles_counts = False
        settings.eval.compare_obstacles_percentage = False #True
        settings.eval.compare_attributes = False#True

        logger.info("The plots are being built now!")
        comparator.process_all_outputs()
        logger.info(
            f"Comparison finished! You will find the results in the 'scenario_comparison' folder"
        )
        end_time = time.time()
        elapsed_time = end_time - start_time
        logger.info(f"Scenario_comparator output processing took {elapsed_time:.2f} seconds")
    
    def add_title_slide(self, prs):
        """
        This method adds a title slide that lists the specific interventions
        or measures being analysed in the subsequent slides, as defined in
        the configuration.

        Parameters
        ----------
        prs : pptx.Presentation
            The presentation object.
        """
        slide = self.duplicate_and_add_slide(prs, 0)
        slide.shapes[0].text_frame.paragraphs[0].runs[0].text = " "
        slide.shapes[1].text_frame.paragraphs[0].runs[0].text = " "
        
        title_box = slide.shapes.add_textbox(Cm(3), Cm(3), Cm(30), Cm(16))
        tf = title_box.text_frame
        # Display intervention
        for intervention in self.configuration.Maßnahmen.split(","):
            para = tf.add_paragraph()
            run = para.add_run()
            run.text = _(intervention.strip())
            run.font.size = Pt(settings.slide.fontsize_pt)
            run.font.bold = settings.slide.fontface_bold
            run.font.name = settings.slide.font


    def add_section_start(self, prs):
        """
        Prepares the initial XML structure for defining presentation sections.
        
        This method is called at the beginning of the presentation creation
        process. It initialises an XML string fragment with the opening tags
        required by the Office Open XML standard for a section list. This
        fragment is then built upon by subsequent calls to `add_section`.

        Parameters
        ----------
        prs : pptx.Presentation
            The presentation object being built.
        """
        xml = '  <p:ext ' + namespacesFragment(["p"])  

        # ext URI has to be {521415D9-36F7-43E2-AB2F-B90AF26B5E84} as it's a registered extension
        xml += '    uri="{521415D9-36F7-43E2-AB2F-B90AF26B5E84}">\n'
        xml += '    <p14:sectionLst ' + namespacesFragment(["p14"]) + '>\n'
        self.section_xml = xml


    def add_section(self, prs):
        """
        Adds a new section divider to the presentation.

        This method uses low-level XML manipulation to insert a section break
        into the PowerPoint file, which helps to organise the presentation.
        The section title is generated based on the configuration.

        Parameters
        ----------
        prs : pptx.Presentation
            The presentation object.
        """
        section_title = (" + ".join([_(variant.strip()) for variant in
                                     self.configuration.Varianten.split(",")]) + ": " if \
                                     settings.slides.section_title_variants else "") + \
                                "_".join(self.configuration.Maßnahmen.split(","))
            
        # section URI's just need to be a GUID wrapped in braces
        xml = (
            '      <p14:section name="'
            + section_title + '" id="{'
            + str(uuid.uuid4()).upper()
            + '}">\n'
        )
        # Only the first slide in the section is added - as section will continue until the next section
        # anyway
        xml += "        <p14:sldIdLst>\n"
        xml += '          <p14:sldId id="' + str(prs.slides[-1].slide_id) + '" />\n'
        xml += "        </p14:sldIdLst>\n"
        xml += "      </p14:section>\n"
        self.section_xml += xml


    def add_section_end(self, prs):
        """
        Finalises and injects the section definitions into the presentation.

        This method is called after all slides and section breaks have been
        added. It closes the XML tags for the section list, creating a
-        complete XML fragment. It then parses this fragment and injects it
        into the main XML of the presentation file, which makes the sections
        appear in the final PowerPoint file.

        Parameters
        ----------
        prs : pptx.Presentation
            The presentation object.
        """
        # Close out the section list
        xml = "    </p14:sectionLst>\n"
    
        # Close out the sections extension
        xml += "  </p:ext>\n"
        self.section_xml += xml
        parsed_xml = parse_xml(self.section_xml)
        
        extLst = OxmlElement("p:extLst")
        prs._element.insert(-1, extLst)

        # Insert the fragment in the extension list in presentation.xml
        extLst.insert(0, parsed_xml)
        

    def move_slide(self, prs, old_index, new_index):
        """
        Moves a slide from one position to another within the presentation.

        Parameters
        ----------
        prs : pptx.Presentation
            The presentation object.
        old_index : int
            The current index of the slide to move.
        new_index : int
            The target index for the slide.
        """
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[old_index])
        xml_slides.insert(new_index, slides[old_index])
        
    def duplicate_and_add_slide(self, pres, index):
        """
        This utility method creates a new, blank slide using the layout of a
        template slide at a given `index`, and then performs a deep copy of all
        shapes from the template onto the new slide. This is used to preserve
        complex slide masters and layouts.

        Parameters
        ----------
        pres : pptx.Presentation
            The presentation object.
        index : int
            The index of the template slide to duplicate.

        Returns
        -------
        pptx.slide.Slide
            The newly created slide object.
        """
        template = pres.slides[index]
    
        # Add a new slide
        copied_slide = pres.slides.add_slide(template.slide_layout)
        
        # Delete the existing shapes that are part of the layout
        for shp in copied_slide.shapes:
            copied_slide.shapes.element.remove(shp.element)        
        
        # Perform a deep copy of the shapes from the template
        for shp in template.shapes:       
            el = shp.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
 
        return copied_slide

    def add_slides(self, prs):
        """
        Adds all specified plots as slides to the presentation.

        This method iterates through the plots listed in the configuration.
        For each plot, it finds the corresponding image file in the output
        folder and adds it to a new slide in the presentation. Each new slide
        is created from a template and populated with a title and subtitle.

        Parameters
        ----------
        prs : pptx.Presentation
            The presentation object to which the slides will be added.
        """
        
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        slide_template = prs.slides[0]
        
        for lang in settings.eval.language:
            subtitle_mapping = dict(zip(self.plots_data["plot_names"], self.plots_data[lang]))
            interventions = self.configuration.Maßnahmen.split(",")
            
            title_text = " + ".join([_(variant.strip()) for variant in self.configuration.Varianten.split(",")]) + ": " + \
                            ((_("Baseline + ") if len(interventions) > 1 else "Baseline") if "Baseline" in interventions else "") + \
                            (_("Maßnahmen ") if len([i for i in interventions if i != "Baseline"]) > 1 else
                             _("Maßnahme ") if len([i for i in interventions if i != "Baseline"]) == 1 else "") + \
                            " + ".join([_(intervention.strip()) for intervention in interventions if intervention != "Baseline"])
    
        # --- MAIN IMG FOLDER: Regular plot slides ---
        for pattern in self.plots_data["plot_names"]:
            
        #for self.plots_data["plot_names"][plot] in self.row[]
            found = False
    
            for filename in os.listdir(self.main_img_folder):
                if filename.endswith(".png") and pattern in filename:
                    file_path = os.path.join(self.main_img_folder, filename)
                    
                    slide = self.duplicate_and_add_slide(prs, 0)
    
                    #slide_title = title_mapping.get(pattern, filename)
                    #title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.05), Inches(8), Inches(1))
                    #tf = title_box.text_frame
                    
                    slide.shapes[0].text_frame.paragraphs[0].runs[0].text = title_text
                    slide.shapes[1].text_frame.paragraphs[0].runs[0].text = \
                        subtitle_mapping.get(pattern, filename).replace("[YEAR]", str(settings.eval.target_year))
    
                    pic_left = Cm(1.33)
                    pic_top = Cm(3.79)
                    height = Cm(13.0)
                    slide.shapes.add_picture(file_path, left=pic_left, top=pic_top, height = height)
                    found = True
    
            if not found:
                print(f"No images found for pattern: {pattern}")
    
        # --- ATTRIBUTE IMG FOLDER: Paired attribute slides ---
        attribute_files = [f for f in os.listdir(self.attribute_img_folder) if f.endswith(".png")]
    
        grouped = defaultdict(dict)
        for fname in attribute_files:
            match = re.match(r"(.*)_((aggregate)|(decomposed))\.png", fname)
            if match:
                base = match.group(1)
                suffix = match.group(2)
                grouped[base][suffix] = os.path.join(self.attribute_img_folder, fname)
    
        for base, versions in grouped.items():
            for suffix in ['aggregate', 'decomposed']:  # consistent order
                if suffix not in versions:
                    continue
        
                file_path = versions[suffix]
                slide = self.duplicate_and_add_slide(prs, 0)
                
                # Determine title based on suffix and language
                subtitle_text = ""
                for lang in settings.eval.language:  # usually only one language
                    try:
                        subtitle_text = self.plots_data.loc[
                            self.plots_data["plot_names"] == suffix.capitalize(), lang
                        ].iloc[0]
                    except IndexError:
                        logger.warning(f"No title found for '{suffix}' in language '{lang}'")
                        subtitle_text = f"{suffix.capitalize()} ({lang})"
        
                # Add title to the slide
                slide.shapes[0].text_frame.paragraphs[0].runs[0].text = title_text
                slide.shapes[1].text_frame.paragraphs[0].runs[0].text = subtitle_text
                
                # Add picture to the slide
                pic_left = Cm(1.33)
                pic_top = Cm(3.79)
                height = Cm(13.0)
                slide.shapes.add_picture(file_path, left=pic_left, top=pic_top, height = height)
        

    def read_input_table(self):
        """
        This helper method parses the Excel configuration file to identify the
        specific combination of scenarios and file prefixes that correspond to
        the presentation variant currently being built. It updates the instance
        attributes accordingly.
        """
        
        file_path = f"{get_output_path(runid=settings.main.run_id, subfolder='slides')}/{settings.slides.config}"
        config_df = pd.read_excel(file_path, sheet_name="Sheet1", skiprows=10)
        list_of_plots_df = pd.read_excel(file_path, sheet_name="Sheet2")
        scenario_dict = {
            "DEZ": "Scenario_mix_pellet_heat_pump",
            "KWN": "Scenario_mix_heat_pump_network_cold",
            "MIX": "Scenario_mix_pellet_heat_pump_network",
            "GPJ": "Scenario_mix_GP_Joule",
        }
        
        filtered_df = config_df[config_df[self.configuration.Variante] == 1]
        
        prefix_list = (filtered_df["Variante"] + "_" + filtered_df["Maßnahme"]).tolist()
        unique_variants = filtered_df["Variante"].unique()
        scenario_list = list({scenario_dict[variant] for variant in unique_variants if variant in scenario_dict})
        
        self.prefixes = prefix_list
        self.scenarios = scenario_list
        self.input_table = config_df
        self.plots_data = list_of_plots_df

config_logging()
Slide_maker().perform_all_processing()
    
    